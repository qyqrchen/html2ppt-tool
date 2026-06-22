import streamlit as st
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt
import io
import time
import os
from playwright.sync_api import sync_playwright

@st.cache_resource
def install_playwright():
    os.system("playwright install chromium")

install_playwright()

def html_to_pptx_with_charts(html_content):
    prs = Presentation()
    
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={'width': 1440, 'height': 2000})
        
        print("\n--- 🚀 通用转换引擎启动 ---")
        page.set_content(html_content, wait_until="networkidle", timeout=30000)
        
        try:
            page.wait_for_function("() => typeof Highcharts !== 'undefined' && Highcharts.charts.length > 0", timeout=8000)
            page.wait_for_timeout(2000)
        except Exception:
            pass

        # 🔥 核心升级1：不再解析原始代码，直接提取运行完 JS、填满表格后的终极 HTML 源码！
        rendered_html = page.content()
        soup = BeautifulSoup(rendered_html, 'html.parser')

        # --- 1. 处理封面页 (兼容旧版 hero 和新版 report-header) ---
        hero = soup.find('div', class_='hero') or soup.find('div', class_='header') or soup.find('div', class_='report-header')
        if hero:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_shape = slide.shapes.title
            subtitle_shape = slide.placeholders[1]
            h1 = hero.find('h1')
            sub = hero.find('div', class_='subtitle')
            title_shape.text = h1.get_text(strip=True) if h1 else "数据分析报告"
            subtitle_shape.text = sub.get_text(strip=True) if sub else ""

        # --- 2. 递归解析逻辑 ---
        def process_element(el, text_frame, slide_title_text):
            if el.name in ['p', 'h4']:
                text = el.get_text(strip=True)
                if text:
                    p = text_frame.add_paragraph()
                    p.text = text
                    p.font.size = Pt(14)
                    if el.name == 'h4':
                        p.font.bold = True
                        
            elif el.name in ['ul', 'ol']:
                for li in el.find_all('li'):
                    p = text_frame.add_paragraph()
                    p.text = li.get_text(strip=True)
                    p.level = 1
                    p.font.size = Pt(14)
                    
            elif el.name == 'table':
                table_slide = prs.slides.add_slide(prs.slide_layouts[5]) 
                table_slide.shapes.title.text = slide_title_text + " - 数据明细"
                rows = el.find_all('tr')
                if not rows: return
                cols_count = max(len(row.find_all(['th', 'td'])) for row in rows)
                
                left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(0.8)
                shape = table_slide.shapes.add_table(len(rows), cols_count, left, top, width, height)
                table = shape.table
                
                for r_idx, row in enumerate(rows):
                    cells = row.find_all(['th', 'td'])
                    for c_idx, cell in enumerate(cells):
                        if c_idx < cols_count:
                            cell_tf = table.cell(r_idx, c_idx).text_frame
                            cell_tf.text = cell.get_text(strip=True)
                            for paragraph in cell_tf.paragraphs:
                                paragraph.font.size = Pt(11)
                                
            elif el.name == 'div':
                classes = el.get('class', [])
                
                # 屏蔽网页控制台和封面，防止乱入正文
                if 'hero' in classes or 'header' in classes or 'controls' in classes or 'global-controls' in classes or 'report-header' in classes:
                    return

                # 图表抓拍
                if 'chart-container' in classes or 'chart-wrapper' in classes:
                    if el.get('id') and str(el.get('id')).startswith('chart'):
                        chart_id = el.get('id')
                    else:
                        child_chart = el.find('div', id=lambda x: x and str(x).startswith('chart'))
                        chart_id = child_chart.get('id') if child_chart else None
                        
                    if chart_id:
                        element_handle = page.locator(f"#{chart_id}")
                        if element_handle.count() > 0:
                            element_handle.scroll_into_view_if_needed()
                            time.sleep(0.5) 
                            image_bytes = element_handle.screenshot()
                            image_stream = io.BytesIO(image_bytes)
                            
                            chart_slide = prs.slides.add_slide(prs.slide_layouts[5])
                            chart_slide.shapes.title.text = slide_title_text + " - 核心图表"
                            chart_slide.shapes.add_picture(image_stream, Inches(1), Inches(1.5), width=Inches(8))
                    return 
                    
                # 🔥 核心升级2：兼容新旧版 KPI 数据卡片
                if 'metric-card' in classes or 'kpi-card' in classes:
                    p = text_frame.add_paragraph()
                    num = el.find(class_=lambda x: x and 'num' in x)
                    label = el.find(class_=lambda x: x and 'label' in x)
                    if num and label:
                        p.text = f"📊 {label.get_text(strip=True)} : {num.get_text(strip=True)}"
                    else:
                        p.text = f"📊 {el.get_text(' : ', strip=True)}"
                    p.font.bold = True
                    p.font.size = Pt(16)
                    return
                
                # 🔥 核心升级3：兼容新版副标题段落
                if 'section-desc' in classes:
                    p = text_frame.add_paragraph()
                    p.text = "📝 " + el.get_text(strip=True)
                    p.font.size = Pt(14)
                    return
                    
                # 🔥 核心升级4：兼容新版的重点项 (优先级条目)
                if 'priority-item' in classes:
                    p = text_frame.add_paragraph()
                    p.text = "🎯 " + el.get_text(" ", strip=True)
                    p.font.size = Pt(14)
                    return
                        
                # 洞察结论 (仅保留旧版兼容，新版 analysis-box 会触发下方的穿透，自动解析出 h4 和 p)
                if 'insight-box' in classes:
                    p = text_frame.add_paragraph()
                    p.text = "💡 洞察结论：\n" + el.get_text(" ", strip=True)
                    p.font.bold = True
                    return
                    
                # 无差别向下钻透
                for child in el.children:
                    if child.name:
                        process_element(child, text_frame, slide_title_text)

        # --- 3. 遍历 Section 生成 ---
        sections = soup.find_all('div', class_='section')
        if not sections:
            container = soup.find('div', class_='container')
            sections = [container] if container else [soup.find('body') or soup]

        for section in sections:
            h_tag = section.find(['h2', 'h1'])
            current_h_text = h_tag.get_text(strip=True) if h_tag else "页面内容提要"
            
            # 过滤掉单纯当做目录的 Section（比如包含 "报告目录" 的 h2）
            if "报告目录" in current_h_text:
                continue
                
            current_slide = prs.slides.add_slide(prs.slide_layouts[1])
            current_slide.shapes.title.text = current_h_text
            text_frame = current_slide.placeholders[1].text_frame
            text_frame.clear()
            
            for el in section.children:
                if el.name is None or el.name in ['h2', 'h1']: 
                    continue
                if el.name == 'h3':
                    current_slide = prs.slides.add_slide(prs.slide_layouts[1])
                    current_slide.shapes.title.text = current_h_text + " \n " + el.get_text(strip=True)
                    current_slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(24)
                    text_frame = current_slide.placeholders[1].text_frame
                    text_frame.clear()
                else:
                    process_element(el, text_frame, current_h_text)
                    
        browser.close()

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output

# --- 🚀 华丽的 Web 界面构建 ---
st.set_page_config(page_title="HTML to PPT 转换器", page_icon="🚀", layout="wide")

with st.sidebar:
    st.image("https://img.icons8.com/color/96/000000/microsoft-powerpoint-2019--v1.png", width=60)
    st.header("⚙️ 转换器引擎状态")
    st.success("✅ 视觉抓拍引擎 (Playwright) 已挂载")
    st.success("✅ DOM无限穿透引擎 已就绪")
    st.markdown("---")
    st.markdown("### 💡 智能识别能力：")
    st.markdown("- **无限穿透**：无视前端外壳嵌套，直击数据底层\n- **文本结构**：各级标题、段落排版\n- **复杂表格**：原生转换为 PPT 表格格式\n- **动态图表**：自动定位、滚动并高清抓拍")
    st.markdown("---")
    st.caption("运行环境：Python 3.12 虚拟隔离舱")

st.title("🚀 HTML to PPT 转换器 (智能容器识别版)")
st.markdown("欢迎使用全自动汇报生成流水线。只需拖入 HTML 文件，系统将自动穿透 DOM 节点，完成动态图表渲染与幻灯片排版。")
st.markdown("---")

uploaded_file = st.file_uploader("📂 请上传需要转换的HTML文件：", type=["html", "htm"])

if uploaded_file is not None:
    original_filename = uploaded_file.name
    output_filename = original_filename.replace(".html", ".pptx").replace(".htm", ".pptx")
    
    html_bytes = uploaded_file.read()
    html_string = html_bytes.decode("utf-8")
    
    st.info(f"📄 文件 `{original_filename}` 载入成功！转换流水线已待命。")
    
    if st.button("⚡ 启动引擎并生成 PPT", type="primary"):
        with st.spinner('底层引擎全速运转中！正在无限穿透嵌套层级，请稍候...'):
            pptx_file = html_to_pptx_with_charts(html_string)
            
        st.balloons()
        st.success(f"🎉 转换圆满完成！已生成：**{output_filename}**")
        
        st.download_button(
            label="⬇️ 点击下载 PPTX 文件",
            data=pptx_file,
            file_name=output_filename,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
